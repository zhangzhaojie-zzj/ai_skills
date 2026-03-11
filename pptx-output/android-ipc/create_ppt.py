#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Android IPC PPT 生成器
使用 python-pptx 创建演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# 配色方案 - Teal Trust 主题
COLORS = {
    'primary': RGBColor(0x02, 0x80, 0x90),      # 深青色
    'secondary': RGBColor(0x00, 0xA8, 0x96),    # 海泡绿
    'accent': RGBColor(0x02, 0xC3, 0x9A),       # 薄荷绿
    'dark': RGBColor(0x1A, 0x1A, 0x2E),         # 深色背景
    'light': RGBColor(0xF8, 0xF9, 0xFA),        # 浅色背景
    'text': RGBColor(0x2C, 0x3E, 0x50),         # 文字颜色
    'text_light': RGBColor(0x5D, 0x6D, 0x7E),   # 次要文字
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'cream': RGBColor(0xFF, 0xF8, 0xE7),
    'light_blue': RGBColor(0xF0, 0xF8, 0xFF),
    'light_green': RGBColor(0xF0, 0xFF, 0xF4),
}

def add_title_slide(prs, title, subtitle, source):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()
    
    # 装饰圆形1
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(4), Inches(4)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = COLORS['secondary']
    circle1.fill.fore_color.brightness = 0.3
    circle1.line.fill.background()
    
    # 装饰圆形2
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent']
    circle2.fill.fore_color.brightness = 0.25
    circle2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # 来源
    source_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = source_box.text_frame
    tf.text = source
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_toc_slide(prs):
    """添加目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "目录"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 目录项
    toc_items = [
        ("01", "跨进程通信方式概览", "Linux IPC 与 Android 特有机制"),
        ("02", "Messenger", "轻量级 IPC 方案"),
        ("03", "AIDL", "接口定义语言实现跨进程"),
        ("04", "ContentProvider", "数据共享的 IPC 方式"),
        ("05", "共享内存", "MemoryFile 与 SharedMemory"),
        ("06", "选型建议", "不同场景的最佳选择"),
    ]
    
    y_pos = 1.3
    for num, title, desc in toc_items:
        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['primary'] if int(num) <= 2 else COLORS['secondary']
        circle.line.fill.background()
        
        # 编号文字
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(0.6), Inches(0.6))
        tf = num_box.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(4), Inches(0.4))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLORS['text']
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos + 0.35), Inches(6), Inches(0.3))
        tf = desc_box.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_light']
        
        y_pos += 0.85
    
    return slide

def add_overview_slide(prs):
    """添加 IPC 概览页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "跨进程通信方式概览"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Linux IPC 卡片
    linux_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5)
    )
    linux_card.fill.solid()
    linux_card.fill.fore_color.rgb = COLORS['white']
    linux_card.line.color.rgb = COLORS['secondary']
    linux_card.line.width = Pt(2)
    
    linux_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.5))
    tf = linux_title.text_frame
    tf.text = "Linux IPC 机制"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    linux_items = [
        "管道 - 缓冲区有限",
        "信号 - 适用于中断控制",
        "信号量 - 同步手段",
        "共享内存 - 速度快",
        "消息队列 - CPU 消耗大",
        "套接字 - 通用但效率低"
    ]
    
    y_pos = 2.1
    for item in linux_items:
        item_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + item
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        y_pos += 0.65
    
    # Android IPC 卡片
    android_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5)
    )
    android_card.fill.solid()
    android_card.fill.fore_color.rgb = COLORS['white']
    android_card.line.color.rgb = COLORS['accent']
    android_card.line.width = Pt(2)
    
    android_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.6), Inches(0.5))
    tf = android_title.text_frame
    tf.text = "Android 特有 IPC"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    android_items = [
        "Messenger - 基于 AIDL",
        "AIDL - 接口定义语言",
        "ContentProvider - 数据共享",
        "MemoryFile - 共享内存",
        "SharedMemory - Android 8+",
        "Binder - 底层实现"
    ]
    
    y_pos = 2.1
    for item in android_items:
        item_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + item
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        y_pos += 0.65
    
    return slide

def add_messenger_slide(prs):
    """添加 Messenger 页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "Messenger - 轻量级 IPC"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(0.4))
    tf = subtitle_box.text_frame
    tf.text = "基于 AIDL 实现，通过 Message 对象在不同进程间传递数据"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_light']
    
    # 左侧特点卡片
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6), Inches(5)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    left_card.line.color.rgb = COLORS['accent']
    left_card.line.width = Pt(2)
    
    left_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = left_title.text_frame
    tf.text = "核心特点"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    features = [
        "轻量级 IPC 方案",
        "底层基于 AIDL",
        "使用 Message 传递数据",
        "支持双向通信",
        "适合简单数据交换"
    ]
    
    y_pos = 2.5
    for feature in features:
        item_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + feature
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        y_pos += 0.7
    
    # 右侧组件卡片
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(6), Inches(5)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    right_card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    right_card.line.width = Pt(1)
    
    right_title = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = right_title.text_frame
    tf.text = "关键组件"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']
    
    components = [
        ("ServiceMessenger", "用于向服务端发送消息"),
        ("ClientMessenger", "用于接收服务端回复"),
        ("ClientHandler", "处理服务端返回的消息")
    ]
    
    y_pos = 2.5
    for name, desc in components:
        name_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.4), Inches(0.35))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Consolas"
        p.font.color.rgb = COLORS['primary']
        
        desc_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos + 0.35), Inches(5.4), Inches(0.3))
        tf = desc_box.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_light']
        
        y_pos += 0.9
    
    return slide

def add_aidl_slide(prs):
    """添加 AIDL 页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "AIDL - Android 接口定义语言"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 定义框
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(1)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = COLORS['cream']
    def_box.line.color.rgb = RGBColor(0xFF, 0xC1, 0x07)
    def_box.line.width = Pt(2)
    
    def_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.933), Inches(0.6))
    tf = def_text.text_frame
    tf.text = "AIDL 允许定义客户端和服务端都认可的编程接口，实现跨进程通信"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text']
    p.alignment = PP_ALIGN.CENTER
    
    # 使用步骤标题
    steps_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.333), Inches(0.5))
    tf = steps_title.text_frame
    tf.text = "使用步骤"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 步骤
    steps = [
        ("1", "创建 .aidl 文件", "定义接口方法"),
        ("2", "实现接口", "在服务端实现 Stub"),
        ("3", "暴露接口", "通过 Service 返回 IBinder"),
        ("4", "客户端调用", "绑定服务后调用方法")
    ]
    
    x_pos = 0.5
    for num, title, desc in steps:
        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x_pos + 0.5), Inches(3.1), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['secondary']
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(x_pos + 0.5), Inches(3.1), Inches(0.8), Inches(0.8))
        tf = num_box.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4), Inches(2.8), Inches(0.4))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.4), Inches(2.8), Inches(0.4))
        tf = desc_box.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += 3.1
    
    return slide

def add_contentprovider_slide(prs):
    """添加 ContentProvider 页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "ContentProvider - 数据共享"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(0.4))
    tf = subtitle_box.text_frame
    tf.text = "底层基于 Binder，用于在不同应用间共享数据"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_light']
    
    # 左侧：数据访问方式
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6), Inches(5)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLORS['white']
    left_card.line.color.rgb = COLORS['secondary']
    left_card.line.width = Pt(2)
    
    left_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = left_title.text_frame
    tf.text = "数据访问方式"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    
    methods = [
        "query() - 查询数据",
        "insert() - 插入数据",
        "update() - 更新数据",
        "delete() - 删除数据",
        "getType() - 返回 MIME 类型"
    ]
    
    y_pos = 2.5
    for method in methods:
        item_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + method
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.name = "Consolas"
        p.font.color.rgb = COLORS['text']
        y_pos += 0.7
    
    # 右侧：使用场景
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(6), Inches(5)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLORS['white']
    right_card.line.color.rgb = COLORS['accent']
    right_card.line.width = Pt(2)
    
    right_title = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = right_title.text_frame
    tf.text = "典型应用场景"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    scenarios = [
        "应用间共享联系人",
        "媒体库访问",
        "日历数据同步",
        "跨应用数据查询",
        "系统级数据访问"
    ]
    
    y_pos = 2.5
    for scenario in scenarios:
        item_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + scenario
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        y_pos += 0.7
    
    return slide

def add_sharedmemory_slide(prs):
    """添加共享内存页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "共享内存 - 高性能 IPC"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # MemoryFile 卡片
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLORS['light_blue']
    left_card.line.color.rgb = COLORS['primary']
    left_card.line.width = Pt(2)
    
    left_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.5))
    tf = left_title.text_frame
    tf.text = "MemoryFile"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    left_subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(5.6), Inches(0.35))
    tf = left_subtitle.text_frame
    tf.text = "传统共享内存实现"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_light']
    
    memoryfile_features = [
        "基于 ashmem 驱动",
        "支持文件描述符传递",
        "需要 ParcelFileDescriptor",
        "适用于大文件传输",
        "需要手动管理内存"
    ]
    
    y_pos = 2.4
    for feature in memoryfile_features:
        item_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + feature
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS['text']
        y_pos += 0.7
    
    # SharedMemory 卡片
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLORS['light_green']
    right_card.line.color.rgb = COLORS['accent']
    right_card.line.width = Pt(2)
    
    right_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.6), Inches(0.5))
    tf = right_title.text_frame
    tf.text = "SharedMemory"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    right_subtitle = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.6), Inches(0.35))
    tf = right_subtitle.text_frame
    tf.text = "Android 8.0+ 新 API"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_light']
    
    sharedmemory_features = [
        "更现代的 API 设计",
        "自动内存管理",
        "支持设置保护标志",
        "更简单的使用方式",
        "推荐用于新开发"
    ]
    
    y_pos = 2.4
    for feature in sharedmemory_features:
        item_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.4), Inches(0.5))
        tf = item_box.text_frame
        tf.text = "• " + feature
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS['text']
        y_pos += 0.7
    
    return slide

def add_recommendation_slide(prs):
    """添加选型建议页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "IPC 机制选型建议"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 选型表格
    recommendations = [
        ("简单数据传递", "Messenger", "轻量、易用"),
        ("复杂接口调用", "AIDL", "灵活、功能强"),
        ("数据共享", "ContentProvider", "标准、安全"),
        ("大文件传输", "SharedMemory", "高性能"),
        ("频繁通信", "Binder", "底层、高效")
    ]
    
    y_pos = 1.2
    colors_bg = [RGBColor(0x1E, 0x3A, 0x4C), RGBColor(0x26, 0x46, 0x53)]
    
    for i, (scenario, solution, reason) in enumerate(recommendations):
        # 背景条
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.9)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = colors_bg[i % 2]
        row_bg.line.fill.background()
        
        # 场景
        scenario_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(3), Inches(0.9))
        tf = scenario_box.text_frame
        tf.text = scenario
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        
        # 推荐方案标签
        solution_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(y_pos + 0.15), Inches(3), Inches(0.6)
        )
        solution_bg.fill.solid()
        solution_bg.fill.fore_color.rgb = COLORS['accent']
        solution_bg.line.fill.background()
        
        solution_box = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 0.15), Inches(3), Inches(0.6))
        tf = solution_box.text_frame
        tf.text = solution
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 原因
        reason_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.5), Inches(0.9))
        tf = reason_box.text_frame
        tf.text = reason
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        
        y_pos += 1.0
    
    return slide

def add_summary_slide(prs):
    """添加总结页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "总结"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 要点
    summary_points = [
        "Android IPC 基于 Linux IPC，但提供了更高效的 Binder 机制",
        "Messenger 适合简单场景，AIDL 适合复杂接口",
        "ContentProvider 是数据共享的标准方式",
        "共享内存适合大文件传输，性能最优",
        "根据场景选择合适的 IPC 机制"
    ]
    
    y_pos = 1.4
    for i, point in enumerate(summary_points):
        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['secondary']
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(0.5), Inches(0.5))
        tf = num_box.text_frame
        tf.text = str(i + 1)
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(11), Inches(0.5))
        tf = content_box.text_frame
        tf.text = point
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        
        y_pos += 0.9
    
    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(6), Inches(11.333), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()
    
    # 底部感谢
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5))
    tf = thanks_box.text_frame
    tf.text = "感谢阅读 | 来源：微信公众号「小墙程序员」"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== 创建所有幻灯片 ====================

# 1. 标题页
add_title_slide(prs, 
    "Android 跨进程通信", 
    "(IPC) 机制详解",
    "来源：微信公众号「小墙程序员」| 掘金")

# 2. 目录
add_toc_slide(prs)

# 3. IPC 概览
add_overview_slide(prs)

# 4. Messenger
add_messenger_slide(prs)

# 5. AIDL
add_aidl_slide(prs)

# 6. ContentProvider
add_contentprovider_slide(prs)

# 7. 共享内存
add_sharedmemory_slide(prs)

# 8. 选型建议
add_recommendation_slide(prs)

# 9. 总结
add_summary_slide(prs)

# 保存文件
output_path = "/Users/mubeichen/blog_ai/pptx-output/android-ipc/Android-IPC-机制详解.pptx"
prs.save(output_path)

print("✅ PPT 创建成功！")
print(f"📁 文件位置: {output_path}")
print(f"📊 共 {len(prs.slides)} 页幻灯片")
